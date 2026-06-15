"""Generation des livrables finaux : PDF rapport et presentation PPTX."""

from __future__ import annotations

import re
import sys
from pathlib import Path

PROJECT_ROOT = Path(__file__).resolve().parents[1]
REPORTS_DIR = PROJECT_ROOT / "reports"
FIGURES_DIR = REPORTS_DIR / "figures"


def _markdown_to_html(markdown_text: str) -> str:
    import base64
    import markdown

    def _embed_image(match: re.Match[str]) -> str:
        alt = match.group("alt")
        rel_path = match.group("path").strip()
        if rel_path.startswith(("http://", "https://", "data:")):
            return match.group(0)
        image_path = (REPORTS_DIR / rel_path).resolve()
        if not image_path.is_file():
            return f"<p><em>Figure introuvable : {rel_path}</em></p>"
        encoded = base64.b64encode(image_path.read_bytes()).decode("ascii")
        ext = image_path.suffix.lower().lstrip(".")
        mime = {"png": "image/png", "jpg": "image/jpeg", "jpeg": "image/jpeg"}.get(ext, "image/png")
        return (
            f'<figure style="margin: 1.2rem 0;">'
            f'<img src="data:{mime};base64,{encoded}" alt="{alt}" style="max-width:100%;">'
            f'<figcaption style="color:#475467;font-size:0.9rem;margin-top:0.4rem;">{alt}</figcaption>'
            f"</figure>"
        )

    pattern = re.compile(r"!\[(?P<alt>[^\]]*)\]\((?P<path>[^)\s]+)\)")
    with_images = pattern.sub(_embed_image, markdown_text)
    body = markdown.markdown(
        with_images,
        extensions=["tables", "fenced_code", "nl2br"],
    )
    return f"""<!DOCTYPE html>
<html lang="fr">
<head>
<meta charset="utf-8">
<title>Rapport d'analyse et d'interpretation</title>
<style>
@page {{ size: A4; margin: 2cm; }}
body {{
  font-family: Helvetica, Arial, sans-serif;
  color: #101828;
  line-height: 1.45;
  font-size: 11pt;
}}
h1 {{ color: #157f7f; font-size: 22pt; border-bottom: 2px solid #157f7f; padding-bottom: 0.3rem; }}
h2 {{ color: #157f7f; font-size: 16pt; margin-top: 1.4rem; }}
h3 {{ color: #344054; font-size: 13pt; }}
h4 {{ color: #475467; font-size: 11.5pt; }}
table {{ border-collapse: collapse; width: 100%; margin: 0.8rem 0; font-size: 10pt; }}
th, td {{ border: 1px solid #d0d5dd; padding: 0.35rem 0.5rem; text-align: left; }}
th {{ background: #f2f4f7; }}
code, pre {{ background: #f9fafb; border: 1px solid #eaecf0; border-radius: 4px; }}
pre {{ padding: 0.6rem; white-space: pre-wrap; font-size: 9pt; }}
hr {{ border: none; border-top: 1px solid #d0d5dd; margin: 1.2rem 0; }}
</style>
</head>
<body>{body}</body>
</html>"""


def generate_pdf() -> Path:
    from xhtml2pdf import pisa

    source = REPORTS_DIR / "rapport_technique.md"
    output = REPORTS_DIR / "rapport_technique.pdf"
    html = _markdown_to_html(source.read_text(encoding="utf-8"))
    with output.open("wb") as pdf_file:
        status = pisa.CreatePDF(html, dest=pdf_file, encoding="utf-8")
    if status.err:
        raise RuntimeError(f"Echec generation PDF ({status.err} erreurs)")
    return output


def _add_bullet_slide(slide, title: str, message: str, bullets: list[str], figure: Path | None = None) -> None:
    from pptx.util import Inches, Pt

    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    p = body.paragraphs[0]
    p.text = message
    p.font.bold = True
    p.font.size = Pt(16)
    for item in bullets:
        para = body.add_paragraph()
        para.text = item
        para.level = 0
        para.font.size = Pt(14)
    if figure and figure.exists():
        slide.shapes.add_picture(str(figure), Inches(6.8), Inches(1.4), width=Inches(2.8))


def generate_pptx() -> Path:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    output = REPORTS_DIR / "presentation_finale.pptx"
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    title_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]

    slides_data: list[tuple] = [
        ("title", "Projet Machine Learning M2 CDSD", [
            "Detection de fraude bancaire",
            "Segmentation client",
            "Industrialisation : API + Dashboard",
            "Kokou Godwin TCHAKPANA — M2 CDSD",
        ]),
        ("content", "Contexte et objectifs", "Reduire le risque de fraude et mieux cibler les clients.", [
            "Detection automatique de transactions frauduleuses",
            "Segmentation intelligente des clients",
            "Deploiement via dashboard Streamlit et API FastAPI",
        ], None),
        ("content", "Donnees et contraintes", "Deux jeux de donnees, deux logiques differentes.", [
            "Fraude : ~1 048 575 transactions, taux ~0,11 %",
            "Segmentation : 2 240 clients, profils heterogenes",
            "Fraude = classe rare ; clustering = interpretation metier",
        ], None),
        ("content", "EDA fraude", "La fraude est rare et concentree sur certains types.", [
            "Fraudes surtout sur TRANSFER et CASH_OUT",
            "Incoherences de soldes tres informatives",
            "Accuracy insuffisante comme metrique decisionnelle",
        ], FIGURES_DIR / "01_fraud_rate_by_type.png"),
        ("content", "EDA segmentation", "Les clients se differencient par valeur, canaux et campagnes.", [
            "Revenu, depenses, canaux web/magasin/promo",
            "24 revenus manquants sur Income",
            "Variables constantes exclues (Z_CostContact, Z_Revenue)",
        ], None),
        ("content", "Pipeline ML fraude", "Feature engineering + seuil calibré sur validation.", [
            "Features soldes : origin_error, dest_error, ratios",
            "Split stratifie + validation temporelle (step)",
            "5 modeles : LogReg, RF, XGBoost, LightGBM, MLP",
            "Seuil calibré avec contrainte de rappel minimal",
        ], None),
        ("content", "Resultats fraude", "XGBoost retenu : meilleure PR-AUC (0,990).", [
            "Logistic Regression : PR-AUC 0,767",
            "Random Forest : PR-AUC 0,987",
            "XGBoost : recall 98,3 %, precision 100 %, seuil 0,67",
            "3 FN / 0 FP — analyse SHAP + chiffrage couts FP/FN",
        ], FIGURES_DIR / "02_fraud_model_comparison.png"),
        ("content", "Pipeline clustering", "Segments lisibles plutot que score mathematique pur.", [
            "Imputation Income + encodage categories",
            "K-Means, Agglomerative, GMM, DBSCAN (k=3 a 6)",
            "Modele retenu : GMM k=4 ; DBSCAN ecarte (91 % bruit)",
        ], FIGURES_DIR / "07_clustering_elbow.png"),
        ("content", "Resultats segmentation", "Quatre segments exploitables malgre silhouette moderee.", [
            "Silhouette : 0,21 ; Davies-Bouldin : 2,42",
            "Projection PCA pour visualiser les groupes",
            "Profils contrastes et actionnables",
        ], FIGURES_DIR / "05_cluster_pca_projection.png"),
        ("content", "Profils clients", "Chaque cluster = une action marketing differente.", [
            "Cluster 0 : dormants (1 132) — reactivation",
            "Cluster 1 : promo digitaux (42) — coupons web + A/B test",
            "Cluster 2 : premium reactifs (175) — VIP",
            "Cluster 3 : fideles stables (891) — cross-sell",
        ], FIGURES_DIR / "06_cluster_profile_heatmap.png"),
        ("content", "Industrialisation MVP", "Projet structure pour une demo bout-en-bout.", [
            "Scripts src/models + modeles joblib",
            "API : /predict/fraud et /predict/segment",
            "Dashboard multi-pages + tests pytest",
            "Docker + CI GitHub Actions + monitoring",
        ], None),
        ("content", "Limites et perspectives", "Resultats solides, validation production a poursuivre.", [
            "Dataset potentiellement tres separable — prudence en prod",
            "Validation temporelle complementaire realisee",
            "Protocole campagne cluster 1 documente",
            "Suivi derive et recalibrage seuil avec le metier",
        ], FIGURES_DIR / "09_fraud_shap_summary.png"),
    ]

    for item in slides_data:
        if item[0] == "title":
            _, title, lines = item
            slide = prs.slides.add_slide(title_layout)
            slide.shapes.title.text = title
            subtitle = slide.placeholders[1].text_frame
            subtitle.clear()
            for i, line in enumerate(lines):
                p = subtitle.paragraphs[0] if i == 0 else subtitle.add_paragraph()
                p.text = line
                p.font.size = Pt(18)
            continue

        _, title, message, bullets, figure = item
        slide = prs.slides.add_slide(content_layout)
        _add_bullet_slide(slide, title, message, bullets, figure)

    prs.save(output)
    return output


def main() -> None:
    pdf_path = generate_pdf()
    pptx_path = generate_pptx()
    print(f"PDF genere : {pdf_path}")
    print(f"PPTX genere : {pptx_path}")


if __name__ == "__main__":
    main()
